pip install pyautogui

pip install python-pptx

from pptx import Presentation
from pptx.util import Inches


import pyautogui
import time

# Open Microsoft PowerPoint
pyautogui.hotkey('winleft', 's')
pyautogui.write('PowerPoint')
pyautogui.press('enter')

time.sleep(3)  # Wait for PowerPoint to open (adjust as needed)

# Create a new presentation
pyautogui.hotkey('ctrl', 'n')

# Add a title slide
pyautogui.hotkey('alt', 'h')  # Go to Home tab
pyautogui.press('m')  # Choose Title Slide layout

# Add title and subtitle
pyautogui.press('tab')  # Move to title placeholder
pyautogui.write("Medical Chatbot for Disease Prediction and Diagnosis")
pyautogui.press('tab')  # Move to subtitle placeholder
pyautogui.write("Integrating KNN and Decision Tree Algorithms")

# Save the presentation
pyautogui.hotkey('ctrl', 's')
pyautogui.write("Medical_Chatbot_Presentation")
pyautogui.press('enter')

# Close PowerPoint
pyautogui.hotkey('alt', 'f')
pyautogui.press('x')

print("Presentation created successfully.")


# Create a presentation object
presentation = Presentation()

# Slide 1: Title
slide1 = presentation.slides.add_slide(presentation.slide_layouts[0])  # 0 corresponds to the title slide layout
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "Medical Chatbot for Disease Prediction and Diagnosis"
subtitle.text = "Integrating KNN and Decision Tree Algorithms"

# Slide 2: Introduction
slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])  # 1 corresponds to the title and content layout
title = slide2.shapes.title
content = slide2.placeholders[1]
title.text = "Introduction"
content.text = "Today, I am excited to share with you our project on the development of a Medical Chatbot designed to predict and diagnose diseases efficiently. We aim to improve healthcare outcomes through early disease prediction."

# Slide 3: Project Overview
slide3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide3.shapes.title
content = slide3.placeholders[1]
title.text = "Project Overview"
content.text = "Our Medical Chatbot predicts diseases by asking users relevant questions and considering reported symptoms. We've integrated KNN and Decision Tree algorithms to achieve a commendable accuracy rate of 90%."

# Slide 4: Current Status
slide4 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide4.shapes.title
content = slide4.placeholders[1]
title.text = "Current Status"
content.text = "Our current model demonstrates a robust performance, achieving a remarkable accuracy of 90% as depicted in the graph. Users can input symptoms and answer questions, leading to accurate disease predictions."

# Slide 5: Future Scope
slide5 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide5.shapes.title
content = slide5.placeholders[1]
title.text = "Future Scope"
content.text = "Looking ahead, we plan to enhance our chatbot using advanced deep learning techniques, such as CNN and Image Segmentation. Our vision includes developing an interactive interface to make the chatbot user-friendly and accessible."

# Slide 6: Interface and Website
slide6 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide6.shapes.title
content = slide6.placeholders[1]
title.text = "Interface and Website"
content.text = "We're working towards creating an intuitive interface that allows users to interact seamlessly with the chatbot. The project includes the development of a dedicated website, ensuring easy access for users."

# Slide 7: Image Scanning Feature
slide7 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide7.shapes.title
content = slide7.placeholders[1]
title.text = "Image Scanning Feature"
content.text = "In the future, our chatbot will be equipped with the capability to scan medical reports from images, making the process even more convenient for users. By leveraging advanced technologies, our chatbot will analyze images to provide accurate predictions and medical assistance."

# Slide 8: Benefits
slide8 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide8.shapes.title
content = slide8.placeholders[1]
title.text = "Benefits"
content.text = "Early Disease Detection: Our chatbot facilitates early detection, leading to timely medical intervention.\nConvenient and Quick Predictions: Users can quickly receive predictions by inputting symptoms and interacting with the chatbot.\nAccess to Medical Information and Assistance: The chatbot provides valuable information about diseases and offers medical assistance."

# Slide 9: Challenges
slide9 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide9.shapes.title
content = slide9.placeholders[1]
title.text = "Challenges"
content.text = "Ensuring Accuracy with Diverse Medical Reports: We recognize the challenge of handling diverse medical reports and are actively working to ensure our chatbot maintains accuracy across various scenarios.\nHandling Ethical Considerations and Data Privacy: As we move forward, ethical considerations and data privacy remain top priorities, and we are committed to implementing robust measures to address these concerns."

# Slide 10: Conclusion
slide10 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide10.shapes.title
content = slide10.placeholders[1]
title.text = "Conclusion"
content.text = "In conclusion, our Medical Chatbot represents a significant step towards revolutionizing disease prediction and diagnosis. The integration of advanced algorithms and upcoming features positions our project at the forefront of healthcare innovation."

# Save the presentation
presentation.save("Medical_Chatbot_Presentation.pptx")
